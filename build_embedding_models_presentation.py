"""
Build a PowerPoint presentation comparing pathology embedding foundation models.

Slides:
  1.  Title
  2.  Why the Embedding Model Matters
  3.  Google Path Foundation  (current)
  4.  CONCH / CONCHv1.5       (recommended #1)
  5.  UNI / UNI2              (recommended #2)
  6.  Virchow2                (recommended #3)
  7.  Prov-GigaPath
  8.  PLIP
  9.  Data-Type Compatibility Matrix
  10. Recommendation & Next Steps

Usage:
    python build_embedding_models_presentation.py
"""

import io
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette (matches build_presentation.py) ────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x5C)
TEAL   = RGBColor(0x00, 0x7A, 0x8A)
ORANGE = RGBColor(0xE8, 0x6C, 0x2C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
DGRAY  = RGBColor(0x3A, 0x3A, 0x3A)
SUBHDR = RGBColor(0xBB, 0xD5, 0xE8)
SUBBDY = RGBColor(0xCC, 0xEE, 0xF5)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xC6, 0x28, 0x28)
AMBER  = RGBColor(0xF5, 0x7F, 0x17)

SW = Inches(13.33)
SH = Inches(7.5)
HEADER_H = Inches(1.15)

OUT_PATH = "results/embedding_models_comparison.pptx"


# ── Low-level helpers (identical API to build_presentation.py) ─────────────────

def blank_slide(prs, bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide


def add_rect(slide, x, y, w, h, fill_color, line=False):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line:
        shp.line.color.rgb = fill_color
    else:
        shp.line.fill.background()
    return shp


def add_textbox(slide, text, x, y, w, h,
                size=14, bold=False, italic=False,
                color=DGRAY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_header(slide, title, subtitle=None):
    bar = add_rect(slide, Inches(0), Inches(0), SW, HEADER_H, NAVY)
    tf = bar.text_frame
    tf.word_wrap = False
    tf.margin_top = Pt(10)
    tf.margin_left = Inches(0.35)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = SUBHDR
        r2.font.name = "Calibri"


def add_bullet_block(slide, items, x, y, w, h, size=13):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            header, subs = item
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(6)
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = f"▶  {header}"
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = NAVY
            r.font.name = "Calibri"
            for sub in subs:
                ps = tf.add_paragraph()
                ps.alignment = PP_ALIGN.LEFT
                ps.space_before = Pt(2)
                rs = ps.add_run()
                rs.text = f"      {sub}"
                rs.font.size = Pt(size - 1)
                rs.font.color.rgb = DGRAY
                rs.font.name = "Calibri"
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(4)
            r = p.add_run()
            r.text = f"•  {item}"
            r.font.size = Pt(size)
            r.font.color.rgb = DGRAY
            r.font.name = "Calibri"


def add_tag(slide, text, x, y, color):
    """Small coloured badge label."""
    box = add_rect(slide, x, y, Inches(1.55), Inches(0.32), color)
    tf = box.text_frame
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"


def add_section_card(slide, title, items, x, y, w, h, header_color=TEAL):
    """Titled card with bullet points inside."""
    add_rect(slide, x, y, w, Inches(0.38), header_color)
    add_textbox(slide, title, x + Inches(0.1), y + Pt(4), w - Inches(0.2), Inches(0.35),
                size=11, bold=True, color=WHITE)
    add_rect(slide, x, y + Inches(0.38), w, h - Inches(0.38), LGRAY)
    tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.45),
                                   w - Inches(0.2), h - Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(10)
        r.font.color.rgb = DGRAY
        r.font.name = "Calibri"


# ── Chart helpers ──────────────────────────────────────────────────────────────

def make_compatibility_matrix_png():
    models = ["Path\nFoundation\n(current)", "CONCH\nv1.5", "UNI\n/ UNI2", "Virchow2", "Prov-\nGigaPath", "PLIP"]
    data_types = ["FFPE H&E", "Frozen H&E", "IHC", "Special stains", "Non-human tissue"]

    # 0 = no, 1 = partial, 2 = yes
    matrix = np.array([
        [2, 2, 1, 1, 1],   # Path Foundation
        [2, 2, 2, 2, 1],   # CONCH v1.5
        [2, 1, 2, 1, 1],   # UNI / UNI2
        [2, 1, 1, 1, 1],   # Virchow2
        [2, 1, 1, 1, 0],   # Prov-GigaPath
        [2, 1, 1, 1, 1],   # PLIP
    ])

    colors = {0: "#C62828", 1: "#F57F17", 2: "#2E7D32"}
    labels = {0: "No", 1: "Partial", 2: "Yes"}

    fig, ax = plt.subplots(figsize=(10, 4.2), facecolor="white")
    ax.set_facecolor("white")

    for i, model in enumerate(models):
        for j, dtype in enumerate(data_types):
            val = matrix[i, j]
            rect = plt.Rectangle([j - 0.45, i - 0.4], 0.9, 0.8,
                                  color=colors[val], alpha=0.85)
            ax.add_patch(rect)
            ax.text(j, i, labels[val], ha="center", va="center",
                    fontsize=9, fontweight="bold", color="white")

    ax.set_xlim(-0.6, len(data_types) - 0.4)
    ax.set_ylim(-0.6, len(models) - 0.4)
    ax.set_xticks(range(len(data_types)))
    ax.set_xticklabels(data_types, fontsize=10)
    ax.set_yticks(range(len(models)))
    ax.set_yticklabels(models, fontsize=9)
    ax.xaxis.set_label_position("top")
    ax.xaxis.tick_top()
    ax.tick_params(length=0)
    ax.spines[:].set_visible(False)

    legend_patches = [
        mpatches.Patch(color="#2E7D32", label="Well supported"),
        mpatches.Patch(color="#F57F17", label="Partial / limited"),
        mpatches.Patch(color="#C62828", label="Not supported"),
    ]
    ax.legend(handles=legend_patches, loc="lower right", fontsize=9,
              framealpha=0.9, bbox_to_anchor=(1.0, -0.02))

    ax.set_title("Data-Type Compatibility by Model", fontsize=12,
                 color="#1B3A5C", pad=14, fontweight="bold")
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf


def make_scale_bar_png():
    models    = ["Path Foundation", "PLIP", "UNI / UNI2", "CONCH v1.5", "Virchow2", "Prov-GigaPath"]
    patches   = [100, 200, 200, 300, 600, 1300]   # millions
    colors    = ["#607D8B", "#00838F", "#1B3A5C", "#00838F", "#E86C2C", "#1B3A5C"]

    fig, ax = plt.subplots(figsize=(7, 3.2), facecolor="white")
    ax.set_facecolor("white")
    bars = ax.barh(models, patches, color=colors, height=0.55)
    ax.set_xlabel("Pretraining patches (millions)", fontsize=10)
    ax.set_title("Pretraining Scale", fontsize=11, color="#1B3A5C", fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(labelsize=9)
    for bar, val in zip(bars, patches):
        ax.text(bar.get_width() + 15, bar.get_y() + bar.get_height() / 2,
                f"{val}M", va="center", fontsize=9, color="#333333")
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf


# ── Individual slides ──────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs, NAVY)
    add_rect(slide, Inches(0), Inches(3.0), SW, Inches(0.07), TEAL)
    add_textbox(slide,
                "Pathology Foundation Models",
                Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.6),
                size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Embedding Model Comparison — Strengths, Weaknesses & Data Compatibility",
                Inches(0.8), Inches(3.15), Inches(11.7), Inches(0.7),
                size=15, color=SUBHDR, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "FFPE H&E  ·  Frozen H&E  ·  IHC  ·  Special Stains",
                Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.5),
                size=13, italic=True, color=RGBColor(0x77, 0x99, 0xBB), align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Models covered: Path Foundation · CONCH · UNI · Virchow2 · Prov-GigaPath · PLIP",
                Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
                size=11, color=RGBColor(0x55, 0x77, 0x99), align=PP_ALIGN.CENTER)


def slide_why_it_matters(prs):
    slide = blank_slide(prs)
    add_header(slide, "Why the Embedding Model Matters",
               "The model determines what features get captured — and what gets lost")

    add_bullet_block(slide, [
        ("Embeddings define your cluster quality",
         ["Two tiles that look similar to a pathologist must be close in embedding space",
          "ImageNet-pretrained models miss histological structure — producing diffuse, overlapping clusters",
          "Pathology-specific models trained on H&E patches capture cell morphology, gland structure, stroma patterns"]),
        ("Stain type & tissue preparation affect the embedding space",
         ["FFPE (fixed) tissue has different morphology to frozen sections — cryoartifacts, ice crystal damage",
          "IHC staining adds chromogen signal that most H&E-only models were not trained to encode",
          "Models trained exclusively on FFPE generalise poorly to frozen tissue"]),
        ("Scale of pretraining correlates with separation quality",
         ["Benchmarks show 30–50% improvement in inter-cluster centroid distance vs ImageNet baselines",
          "Larger pretraining datasets yield more compact, better-separated clusters in t-SNE / UMAP"]),
        ("Access",
         ["All top-performing models are gated on HuggingFace (accept license agreement required)",
          "PyTorch + HuggingFace Transformers — existing TF extractor code needs updating"]),
    ], Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8), size=13)


def slide_model(prs, model_name, subtitle, badge_text, badge_color,
                arch, pretraining, strengths, weaknesses, best_for, access_text, access_color):
    slide = blank_slide(prs)
    add_header(slide, model_name, subtitle)

    # Badge
    add_tag(slide, badge_text, Inches(11.6), Inches(0.35), badge_color)

    # Arch + pretraining strip
    add_rect(slide, Inches(0), HEADER_H, SW, Inches(0.42), LGRAY)
    add_textbox(slide,
                f"Architecture: {arch}     |     Pretraining: {pretraining}     |     Access: {access_text}",
                Inches(0.35), HEADER_H + Pt(4), Inches(12.6), Inches(0.38),
                size=10, italic=True, color=DGRAY)

    card_y = HEADER_H + Inches(0.52)
    card_h = Inches(2.1)
    card_w = Inches(4.1)
    gap    = Inches(0.27)

    add_section_card(slide, "Strengths", strengths,
                     Inches(0.3), card_y, card_w, card_h, TEAL)

    add_section_card(slide, "Weaknesses", weaknesses,
                     Inches(0.3) + card_w + gap, card_y, card_w, card_h, ORANGE)

    add_section_card(slide, "Best for (data type)", best_for,
                     Inches(0.3) + (card_w + gap) * 2, card_y, card_w, card_h, NAVY)

    # Bottom stain compatibility tags
    stain_y = card_y + card_h + Inches(0.25)
    add_textbox(slide, "Stain / prep compatibility:",
                Inches(0.3), stain_y, Inches(2.8), Inches(0.35),
                size=11, bold=True, color=NAVY)


def slide_path_foundation(prs):
    slide = blank_slide(prs)
    add_header(slide, "Google Path Foundation", "Current model in use — 384-dim embeddings")
    add_tag(slide, "CURRENT", Inches(11.6), Inches(0.35), TEAL)

    add_rect(slide, Inches(0), HEADER_H, SW, Inches(0.42), LGRAY)
    add_textbox(slide,
                "Architecture: ViT (TF SavedModel)     |     Pretraining: ~100M pathology patches     |     Access: Gated (HuggingFace)     |     Dim: 384",
                Inches(0.35), HEADER_H + Pt(4), Inches(12.6), Inches(0.38),
                size=10, italic=True, color=DGRAY)

    card_y = HEADER_H + Inches(0.52)
    card_h = Inches(2.15)
    card_w = Inches(4.1)
    gap    = Inches(0.27)

    add_section_card(slide, "Strengths", [
        "Easy to load — TF SavedModel, works with existing extractor code",
        "Good baseline separation on standard H&E patches",
        "Lightweight 384-dim vectors — fast clustering & low memory",
        "Publicly documented by Google",
    ], Inches(0.3), card_y, card_w, card_h, TEAL)

    add_section_card(slide, "Weaknesses", [
        "Trained primarily on FFPE H&E — limited frozen section robustness",
        "Smaller pretraining scale vs newer models (UNI, Virchow2)",
        "Benchmarks show weaker organ/tissue separation than CONCH or UNI",
        "CHIEF-style architecture has known cluster overlap on multi-organ data",
    ], Inches(0.3) + card_w + gap, card_y, card_w, card_h, ORANGE)

    add_section_card(slide, "Best for", [
        "FFPE H&E standard slides",
        "Quick prototyping and baseline runs",
        "Memory-constrained environments (384-dim)",
        "TensorFlow-only setups",
    ], Inches(0.3) + (card_w + gap) * 2, card_y, card_w, card_h, NAVY)

    stain_y = card_y + card_h + Inches(0.28)
    add_textbox(slide, "Stain / prep compatibility:", Inches(0.3), stain_y, Inches(2.8), Inches(0.35),
                size=11, bold=True, color=NAVY)
    tags = [("FFPE H&E ✓", GREEN), ("Frozen H&E ~", AMBER), ("IHC ~", AMBER), ("Special stains ✗", RED)]
    for i, (txt, col) in enumerate(tags):
        add_tag(slide, txt, Inches(3.1) + i * Inches(1.7), stain_y, col)

    add_bullet_block(slide, [
        "Recommended when: you need a fast baseline, are using TensorFlow, or have standard FFPE H&E data only.",
        "Upgrade when: you observe poor cluster separation, have frozen sections or IHC slides.",
    ], Inches(0.3), stain_y + Inches(0.45), Inches(12.5), Inches(1.5), size=12)


def slide_conch(prs):
    slide = blank_slide(prs)
    add_header(slide, "CONCH / CONCHv1.5",
               "CONtrastive learning from Captions for Histopathology — Mahmood Lab, Harvard")
    add_tag(slide, "RECOMMENDED #1", Inches(10.9), Inches(0.35), GREEN)

    add_rect(slide, Inches(0), HEADER_H, SW, Inches(0.42), LGRAY)
    add_textbox(slide,
                "Architecture: Vision-Language ViT     |     Pretraining: 1.17M image-caption pairs, diverse stains     |     Access: Gated (HuggingFace)     |     Dim: 512",
                Inches(0.35), HEADER_H + Pt(4), Inches(12.6), Inches(0.38),
                size=10, italic=True, color=DGRAY)

    card_y = HEADER_H + Inches(0.52)
    card_h = Inches(2.15)
    card_w = Inches(4.1)
    gap    = Inches(0.27)

    add_section_card(slide, "Strengths", [
        "Best generalisation to external/unseen datasets in benchmarks",
        "CONCHv1.5 explicitly pretrained on frozen tissue, IHC, and special stains",
        "Top scores on tissue subtyping: 90.7% NSCLC, 90.2% RCC (zero-shot)",
        "Robust across staining protocols — FFPE, frozen, IHC in one model",
        "Intrinsic scale-invariance for tile-level embeddings",
    ], Inches(0.3), card_y, card_w, card_h, TEAL)

    add_section_card(slide, "Weaknesses", [
        "Vision-language architecture — language bottleneck may compress visual detail",
        "Vision-only models (UNI, Virchow2) sometimes score higher on purely visual tasks",
        "Requires PyTorch — existing TF extractor needs updating",
        "Gated access (license agreement on HuggingFace required)",
    ], Inches(0.3) + card_w + gap, card_y, card_w, card_h, ORANGE)

    add_section_card(slide, "Best for", [
        "Frozen H&E sections (primary strength of v1.5)",
        "Multi-stain datasets (H&E + IHC + special stains)",
        "When generalising to new tissue types not seen during training",
        "Clustering with diverse tissue preparations",
    ], Inches(0.3) + (card_w + gap) * 2, card_y, card_w, card_h, NAVY)

    stain_y = card_y + card_h + Inches(0.28)
    add_textbox(slide, "Stain / prep compatibility:", Inches(0.3), stain_y, Inches(2.8), Inches(0.35),
                size=11, bold=True, color=NAVY)
    tags = [("FFPE H&E ✓", GREEN), ("Frozen H&E ✓", GREEN), ("IHC ✓", GREEN), ("Special stains ✓", GREEN)]
    for i, (txt, col) in enumerate(tags):
        add_tag(slide, txt, Inches(3.1) + i * Inches(1.7), stain_y, col)

    add_bullet_block(slide, [
        "Why #1 for your data: CONCHv1.5 is the only top model explicitly pretrained on frozen tissue.",
        "GitHub: github.com/mahmoodlab/CONCH  ·  HuggingFace: MahmoodLab/conch",
    ], Inches(0.3), stain_y + Inches(0.45), Inches(12.5), Inches(1.5), size=12)


def slide_uni(prs):
    slide = blank_slide(prs)
    add_header(slide, "UNI / UNI2",
               "Universal Pathology Model — Mahmood Lab, Harvard · Nature Medicine 2024")
    add_tag(slide, "RECOMMENDED #2", Inches(10.9), Inches(0.35), TEAL)

    add_rect(slide, Inches(0), HEADER_H, SW, Inches(0.42), LGRAY)
    add_textbox(slide,
                "Architecture: ViT-L (DINOv2)     |     Pretraining: 200M+ patches, 350k+ slides (Mass General Brigham)     |     Access: Gated (HuggingFace)     |     Dim: 1024",
                Inches(0.35), HEADER_H + Pt(4), Inches(12.6), Inches(0.38),
                size=10, italic=True, color=DGRAY)

    card_y = HEADER_H + Inches(0.52)
    card_h = Inches(2.15)
    card_w = Inches(4.1)
    gap    = Inches(0.27)

    add_section_card(slide, "Strengths", [
        "Vision-only — preserves richer visual detail than VL models",
        "Trained on 200M+ patches → very strong in-distribution performance",
        "DINOv2 backbone — excellent self-supervised spatial features",
        "UNI2 adds IHC and broader tissue diversity over original UNI",
        "Strong benchmark results across cancer classification tasks",
    ], Inches(0.3), card_y, card_w, card_h, TEAL)

    add_section_card(slide, "Weaknesses", [
        "Trained primarily on FFPE H&E — frozen tissue not a focus",
        "1024-dim vectors: higher memory & compute cost than CONCH / Path Foundation",
        "Strong in-distribution but shows performance drop on external datasets",
        "Requires PyTorch — existing TF extractor needs updating",
    ], Inches(0.3) + card_w + gap, card_y, card_w, card_h, ORANGE)

    add_section_card(slide, "Best for", [
        "FFPE H&E slides — primary training domain",
        "Large-scale clustering where in-distribution performance is key",
        "Cancer subtyping, grading, and morphology clustering",
        "Datasets similar to academic medical centre pathology",
    ], Inches(0.3) + (card_w + gap) * 2, card_y, card_w, card_h, NAVY)

    stain_y = card_y + card_h + Inches(0.28)
    add_textbox(slide, "Stain / prep compatibility:", Inches(0.3), stain_y, Inches(2.8), Inches(0.35),
                size=11, bold=True, color=NAVY)
    tags = [("FFPE H&E ✓", GREEN), ("Frozen H&E ~", AMBER), ("IHC ~", AMBER), ("Special stains ~", AMBER)]
    for i, (txt, col) in enumerate(tags):
        add_tag(slide, txt, Inches(3.1) + i * Inches(1.7), stain_y, col)

    add_bullet_block(slide, [
        "Same HuggingFace access as CONCH (Mahmood Lab) — request both at the same time.",
        "GitHub: github.com/mahmoodlab/UNI  ·  HuggingFace: MahmoodLab/uni",
    ], Inches(0.3), stain_y + Inches(0.45), Inches(12.5), Inches(1.5), size=12)


def slide_virchow2(prs):
    slide = blank_slide(prs)
    add_header(slide, "Virchow2",
               "Large-scale pathology foundation model — Paige AI")
    add_tag(slide, "RECOMMENDED #3", Inches(10.9), Inches(0.35), ORANGE)

    add_rect(slide, Inches(0), HEADER_H, SW, Inches(0.42), LGRAY)
    add_textbox(slide,
                "Architecture: ViT-H/14 (632M params)     |     Pretraining: 3M+ slides (Paige / Memorial Sloan Kettering)     |     Access: Gated (HuggingFace — Paige AI)     |     Dim: 1280",
                Inches(0.35), HEADER_H + Pt(4), Inches(12.6), Inches(0.38),
                size=10, italic=True, color=DGRAY)

    card_y = HEADER_H + Inches(0.52)
    card_h = Inches(2.15)
    card_w = Inches(4.1)
    gap    = Inches(0.27)

    add_section_card(slide, "Strengths", [
        "Largest vision-only model — 632M parameters (ViT-H/14)",
        "Intrinsic scale-invariance: robust without augmentation tricks",
        "Trained on 3M+ clinical slides from MSK — enormous diversity",
        "Top scores across multiple external benchmarks (MSI prediction, cancer typing)",
        "Excels at capturing subtle morphological differences",
    ], Inches(0.3), card_y, card_w, card_h, TEAL)

    add_section_card(slide, "Weaknesses", [
        "1280-dim vectors — highest memory & compute cost of all models",
        "FFPE-heavy training data — frozen section generalisation not validated",
        "Gated via Paige AI (separate access request from Mahmood Lab models)",
        "Larger model = slower inference per batch",
    ], Inches(0.3) + card_w + gap, card_y, card_w, card_h, ORANGE)

    add_section_card(slide, "Best for", [
        "FFPE H&E — primary training domain",
        "Tasks requiring fine-grained morphological separation",
        "Large compute environments (GPU with >16 GB VRAM recommended)",
        "When CONCH and UNI do not give sufficient cluster separation",
    ], Inches(0.3) + (card_w + gap) * 2, card_y, card_w, card_h, NAVY)

    stain_y = card_y + card_h + Inches(0.28)
    add_textbox(slide, "Stain / prep compatibility:", Inches(0.3), stain_y, Inches(2.8), Inches(0.35),
                size=11, bold=True, color=NAVY)
    tags = [("FFPE H&E ✓", GREEN), ("Frozen H&E ~", AMBER), ("IHC ~", AMBER), ("Special stains ~", AMBER)]
    for i, (txt, col) in enumerate(tags):
        add_tag(slide, txt, Inches(3.1) + i * Inches(1.7), stain_y, col)

    add_bullet_block(slide, [
        "Access: huggingface.co/paige-ai/Virchow2  — requires separate agreement with Paige AI.",
        "Only worth the extra effort if CONCH (#1) and UNI (#2) do not produce satisfying cluster separation.",
    ], Inches(0.3), stain_y + Inches(0.45), Inches(12.5), Inches(1.5), size=12)


def slide_provgigapath(prs):
    slide = blank_slide(prs)
    add_header(slide, "Prov-GigaPath",
               "Whole-slide pathology foundation model — Microsoft / Providence Health")

    add_rect(slide, Inches(0), HEADER_H, SW, Inches(0.42), LGRAY)
    add_textbox(slide,
                "Architecture: ViT-g + LongNet (slide-level)     |     Pretraining: 1.3B images, 170k+ slides (Providence)     |     Access: Gated (HuggingFace)     |     Dim: 1536",
                Inches(0.35), HEADER_H + Pt(4), Inches(12.6), Inches(0.38),
                size=10, italic=True, color=DGRAY)

    card_y = HEADER_H + Inches(0.52)
    card_h = Inches(2.15)
    card_w = Inches(4.1)
    gap    = Inches(0.27)

    add_section_card(slide, "Strengths", [
        "Largest pretraining dataset: 1.3 billion pathology images",
        "Designed for slide-level reasoning — captures global tissue context",
        "SOTA on cancer subtyping when using slide-level aggregation",
        "Backed by Microsoft Research — strong engineering & documentation",
    ], Inches(0.3), card_y, card_w, card_h, TEAL)

    add_section_card(slide, "Weaknesses", [
        "Optimised for slide-level tasks — tile-level embeddings less competitive",
        "Benchmarks show poor organ-level embedding separation at tile level",
        "1536-dim — highest dimensionality of all models listed",
        "Primarily US clinical FFPE data — limited diversity of tissue preparation",
        "Not tested on frozen sections or non-human tissue",
    ], Inches(0.3) + card_w + gap, card_y, card_w, card_h, ORANGE)

    add_section_card(slide, "Best for", [
        "Slide-level classification and survival prediction tasks",
        "FFPE H&E with large-scale clinical datasets",
        "MIL (multiple instance learning) aggregation workflows",
        "NOT recommended for tile-level clustering",
    ], Inches(0.3) + (card_w + gap) * 2, card_y, card_w, card_h, NAVY)

    stain_y = card_y + card_h + Inches(0.28)
    add_textbox(slide, "Stain / prep compatibility:", Inches(0.3), stain_y, Inches(2.8), Inches(0.35),
                size=11, bold=True, color=NAVY)
    tags = [("FFPE H&E ✓", GREEN), ("Frozen H&E ~", AMBER), ("IHC ~", AMBER), ("Special stains ✗", RED)]
    for i, (txt, col) in enumerate(tags):
        add_tag(slide, txt, Inches(3.1) + i * Inches(1.7), stain_y, col)

    add_bullet_block(slide, [
        "Verdict: Not recommended for tile-level clustering. Benchmarks confirm poor tile-level separation despite massive pretraining scale.",
    ], Inches(0.3), stain_y + Inches(0.45), Inches(12.5), Inches(1.0), size=12)


def slide_plip(prs):
    slide = blank_slide(prs)
    add_header(slide, "PLIP",
               "Pathology Language-Image Pretraining — Stanford · Nature Medicine 2023")

    add_rect(slide, Inches(0), HEADER_H, SW, Inches(0.42), LGRAY)
    add_textbox(slide,
                "Architecture: CLIP (ViT-B/32)     |     Pretraining: 200k image-text pairs from Twitter/social media     |     Access: Open (HuggingFace)     |     Dim: 512",
                Inches(0.35), HEADER_H + Pt(4), Inches(12.6), Inches(0.38),
                size=10, italic=True, color=DGRAY)

    card_y = HEADER_H + Inches(0.52)
    card_h = Inches(2.15)
    card_w = Inches(4.1)
    gap    = Inches(0.27)

    add_section_card(slide, "Strengths", [
        "Fully open weights — no gated access required",
        "Vision-language: supports text-guided retrieval and zero-shot classification",
        "Good zero-shot performance on common tissue types",
        "Smallest & fastest model — easy to deploy",
        "Well documented, widely used baseline in research",
    ], Inches(0.3), card_y, card_w, card_h, TEAL)

    add_section_card(slide, "Weaknesses", [
        "Trained on Twitter image-text pairs — noisy, limited tissue diversity",
        "ViT-B/32 backbone — weaker than ViT-L/H models (UNI, Virchow2)",
        "CONCH outperforms PLIP by 12% on NSCLC subtyping, 9.8% on RCC",
        "No specific support for frozen tissue or IHC",
        "Limited to ~200k training pairs vs millions for newer models",
    ], Inches(0.3) + card_w + gap, card_y, card_w, card_h, ORANGE)

    add_section_card(slide, "Best for", [
        "Quick open-source baseline (no HuggingFace gated access needed)",
        "Text-guided image retrieval (finding similar tiles by description)",
        "Standard FFPE H&E datasets",
        "Resource-constrained environments",
    ], Inches(0.3) + (card_w + gap) * 2, card_y, card_w, card_h, NAVY)

    stain_y = card_y + card_h + Inches(0.28)
    add_textbox(slide, "Stain / prep compatibility:", Inches(0.3), stain_y, Inches(2.8), Inches(0.35),
                size=11, bold=True, color=NAVY)
    tags = [("FFPE H&E ✓", GREEN), ("Frozen H&E ~", AMBER), ("IHC ~", AMBER), ("Special stains ~", AMBER)]
    for i, (txt, col) in enumerate(tags):
        add_tag(slide, txt, Inches(3.1) + i * Inches(1.7), stain_y, col)

    add_bullet_block(slide, [
        "Verdict: Best open-source baseline. For clustering quality, all gated models above outperform PLIP significantly.",
    ], Inches(0.3), stain_y + Inches(0.45), Inches(12.5), Inches(1.0), size=12)


def slide_compatibility_matrix(prs):
    slide = blank_slide(prs)
    add_header(slide, "Data-Type Compatibility Matrix",
               "Which model works with which tissue preparation and stain")

    matrix_buf = make_compatibility_matrix_png()
    slide.shapes.add_picture(matrix_buf, Inches(0.3), Inches(1.25),
                              Inches(8.5), Inches(4.2))

    scale_buf = make_scale_bar_png()
    slide.shapes.add_picture(scale_buf, Inches(9.0), Inches(1.25),
                              Inches(4.1), Inches(3.2))

    add_bullet_block(slide, [
        "CONCH v1.5 is the only model explicitly pretrained on frozen tissue — making it uniquely suited to your dataset.",
        "UNI and Virchow2 have partial frozen support but were not optimised for it.",
        "Prov-GigaPath has the largest pretraining scale but the weakest tile-level separation.",
    ], Inches(0.3), Inches(5.7), Inches(12.5), Inches(1.6), size=12)


def slide_recommendation(prs):
    slide = blank_slide(prs)
    add_header(slide, "Recommendation & Next Steps",
               "Frozen H&E clustering — GI frozen sections, 512×512 px tiles")

    steps = [
        ("1", "CONCH v1.5", "Try first\nFrozen tissue optimised", GREEN),
        ("2", "UNI / UNI2", "Try second\nLargest H&E dataset", TEAL),
        ("3", "Virchow2", "Try third\nIf #1 & #2 insufficient", ORANGE),
    ]

    bw, bh = Inches(3.6), Inches(1.8)
    gap = Inches(0.35)
    sx = Inches(0.4)
    for i, (num, name, desc, color) in enumerate(steps):
        x = sx + i * (bw + gap)
        box = add_rect(slide, x, Inches(1.4), bw, bh, color)
        tf = box.text_frame
        tf.margin_top = Pt(12)
        tf.margin_left = Pt(14)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"#{num}  {name}"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "Calibri"
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = SUBBDY
        r2.font.name = "Calibri"

    add_bullet_block(slide, [
        ("Access all three at once",
         ["Request access on HuggingFace for MahmoodLab/conch, MahmoodLab/uni, and paige-ai/Virchow2",
          "Mahmood Lab models share the same approval — request both CONCH and UNI simultaneously"]),
        ("Code change required",
         ["All three use PyTorch + HuggingFace Transformers",
          "The embedding_extractor.py TF/Keras model loader needs to be replaced with a HuggingFace pipeline",
          "The rest of the clustering pipeline (DuckDB, KMeans, CLI) is model-agnostic — no changes needed"]),
        ("Evaluation approach",
         ["Run find_optimal_k_advanced.py with each model's embeddings",
          "Compare CHI, DBI, and silhouette scores across models",
          "Visually inspect cluster_grids for morphological coherence"]),
    ], Inches(0.4), Inches(3.4), Inches(12.5), Inches(3.9), size=12)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building slides …")
    slide_title(prs)                ; print("  1/10  Title")
    slide_why_it_matters(prs)       ; print("  2/10  Why It Matters")
    slide_path_foundation(prs)      ; print("  3/10  Path Foundation (current)")
    slide_conch(prs)                ; print("  4/10  CONCH / CONCHv1.5")
    slide_uni(prs)                  ; print("  5/10  UNI / UNI2")
    slide_virchow2(prs)             ; print("  6/10  Virchow2")
    slide_provgigapath(prs)         ; print("  7/10  Prov-GigaPath")
    slide_plip(prs)                 ; print("  8/10  PLIP")
    slide_compatibility_matrix(prs) ; print("  9/10  Compatibility Matrix")
    slide_recommendation(prs)       ; print("  10/10 Recommendation & Next Steps")

    os.makedirs("results", exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\nSaved: {OUT_PATH}")


if __name__ == "__main__":
    main()
