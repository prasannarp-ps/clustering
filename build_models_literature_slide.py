"""
Generate a PowerPoint presentation with two tables:
  1. Pathology Foundation Models overview
  2. Key Literature References

Usage:
    python build_models_literature_slide.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def set_cell(cell, text, bold=False, size=Pt(10), font_color=None, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    if font_color:
        run.font.color.rgb = font_color


def style_header_row(table, col_count, bg_color=RGBColor(0x1B, 0x5E, 0x20)):
    for i in range(col_count):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        set_cell(cell, cell.text, bold=True, size=Pt(11),
                 font_color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)


def stripe_rows(table, row_count, col_count):
    for r in range(1, row_count):
        color = RGBColor(0xF5, 0xF5, 0xF5) if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        for c in range(col_count):
            table.cell(r, c).fill.solid()
            table.cell(r, c).fill.fore_color.rgb = color


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =========================================================
    # SLIDE 1 — Foundation Models Table
    # =========================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Pathology Foundation Models — Overview"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    models = [
        ("CONCH / v1.5", "ViT-B (CoCa)", "512", "1.17M image-caption pairs\n(PubMed, books, web)",
         "Best generalizability across stains\n(H&E, IHC, frozen). Top benchmark\nscores on external datasets.",
         "github.com/mahmoodlab/CONCH"),
        ("UNI / UNI 2", "ViT-L (DINOv2)", "1024 / 1536", "200M+ patches from\n350k WSIs",
         "Strong in-distribution separation.\nUNI2 is the strongest variant\n(use batch_size 64).",
         "github.com/mahmoodlab/UNI"),
        ("Virchow2", "ViT-H/14\n(632M params)", "1280", "3.1M WSIs\n(Paige clinical archive)",
         "Largest vision-only model.\nIntrinsic scale robustness\nwithout augmentation.",
         "huggingface.co/paige-ai/Virchow2"),
        ("Prov-GigaPath", "ViT-G (DINOv2)", "1536", "1.3B pathology tiles\n(Providence Health)",
         "SOTA on cancer classification.\nSlide-level + tile-level\npretrained jointly.",
         "github.com/prov-gigapath/prov-gigapath"),
        ("Kaiko\n(Midnight)", "ViT-L/G", "1024 / 1536", "Large-scale pathology\ncorpus",
         "Competitive with Virchow2 on\nseveral benchmarks.",
         "github.com/kaiko-ai/Midnight"),
        ("Hibou-L", "ViT-L (DINOv2)", "1024", "Pathology tile corpus", "Performs well on tile-level\ntasks.",
         "github.com/HistAI/hibou"),
        ("PathDino", "ViT-S (DINOv2)", "384", "Pathology fine-tuned\nDINOv2",
         "Lightweight. Good for fast\ntile-level feature extraction.",
         "github.com/KimiaLabMayo/PathDino"),
        ("PLIP", "ViT-B/32 (CLIP)", "512", "Twitter pathology\nimage-caption pairs",
         "Language-aligned embeddings.\nUseful for zero-shot retrieval.",
         "github.com/PathologyFoundation/plip"),
        ("QuiltNet", "ViT-B/32 (CLIP)", "512", "1M YouTube pathology\nvideo-text pairs",
         "Language-aligned. Good for\ncross-modal search tasks.",
         "huggingface.co/wisdomik/QuiltNet-B-32"),
        ("TITAN", "ViT-B (CoCa)\n+ WSI aggregator", "512", "WSI-level multimodal\npretraining",
         "Better for slide-level tasks\nthan tile-level clustering.\nInternal CONCH v1.5 encoder.",
         "github.com/mahmoodlab/TITAN"),
    ]

    headers = ["Model", "Architecture", "Embed\nDim", "Training Data", "Key Strengths", "Repository"]
    col_widths = [Inches(1.4), Inches(1.5), Inches(0.8), Inches(2.2), Inches(3.2), Inches(3.5)]
    n_rows = len(models) + 1
    n_cols = len(headers)

    table_shape = slide1.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(0.9),
                                           sum(col_widths), Inches(6.2))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = int(w)

    for i, h in enumerate(headers):
        set_cell(table.cell(0, i), h, bold=True, size=Pt(11),
                 font_color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

    for r, row_data in enumerate(models, start=1):
        for c, val in enumerate(row_data):
            alignment = PP_ALIGN.CENTER if c == 2 else PP_ALIGN.LEFT
            set_cell(table.cell(r, c), val, size=Pt(9), alignment=alignment)

    style_header_row(table, n_cols)
    stripe_rows(table, n_rows, n_cols)

    # =========================================================
    # SLIDE 2 — Literature References Table
    # =========================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Key Literature — Pathology Foundation Models"
    run2.font.size = Pt(24)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    references = [
        ("A multimodal whole-slide foundation\nmodel for pathology",
         "Chen et al., 2025",
         "Nature Medicine",
         "Introduces TITAN: multimodal WSI-level model combining\nvision encoder with language-guided pretraining.",
         "nature.com/articles/s41591-025-03982-3"),
        ("Pathology Foundation Models",
         "Vrabac et al., 2024",
         "J Pathol Inform",
         "Comprehensive review of foundation models in pathology.\nCovers architectures, training strategies, and benchmarks.",
         "pmc.ncbi.nlm.nih.gov/articles/PMC11799676"),
        ("Comparing Computational Pathology\nFoundation Models using\nRepresentational Similarity Analysis",
         "Unger et al., 2025",
         "arXiv preprint",
         "RSA-based comparison showing which models learn similar\nrepresentations. Useful for understanding model redundancy.",
         "arxiv.org/abs/2509.15482"),
        ("Benchmarking foundation models as\nfeature extractors for weakly-supervised\ncomputational pathology",
         "Fourkioti et al., 2024",
         "arXiv preprint",
         "Systematic benchmark of foundation models as frozen feature\nextractors. Evaluates on multiple downstream tasks.",
         "arxiv.org/abs/2408.15823"),
        ("A whole-slide foundation model for\ndigital pathology from real-world data",
         "Vorontsov et al., 2024",
         "Nature",
         "Introduces Virchow: ViT-H model trained on 1.5M slides.\nDemonstrates scale advantages for pathology FMs.",
         "nature.com/articles/s41586-024-07441-w"),
        ("A visual-language foundation model\nfor computational pathology",
         "Lu et al., 2024",
         "Nature Medicine",
         "Introduces CONCH: CoCa-based model with vision-language\nalignment for zero-shot and few-shot pathology tasks.",
         "pmc.ncbi.nlm.nih.gov/articles/PMC11384335"),
    ]

    ref_headers = ["Title", "Authors", "Venue", "Summary", "Link"]
    ref_col_widths = [Inches(2.8), Inches(1.5), Inches(1.2), Inches(4.5), Inches(2.8)]
    n_ref_rows = len(references) + 1
    n_ref_cols = len(ref_headers)

    table_shape2 = slide2.shapes.add_table(n_ref_rows, n_ref_cols, Inches(0.4), Inches(0.9),
                                            sum(ref_col_widths), Inches(5.5))
    table2 = table_shape2.table

    for i, w in enumerate(ref_col_widths):
        table2.columns[i].width = int(w)

    for i, h in enumerate(ref_headers):
        set_cell(table2.cell(0, i), h, bold=True, size=Pt(11),
                 font_color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

    for r, row_data in enumerate(references, start=1):
        for c, val in enumerate(row_data):
            set_cell(table2.cell(r, c), val, size=Pt(9))

    style_header_row(table2, n_ref_cols, bg_color=RGBColor(0x0D, 0x47, 0xA1))
    stripe_rows(table2, n_ref_rows, n_ref_cols)

    # Save
    out_path = "results/pathology_foundation_models_overview.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    build_presentation()
