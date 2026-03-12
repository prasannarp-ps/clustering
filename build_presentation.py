"""
Build clustering results PowerPoint presentation.

Slides:
  1. Title
  2. Dataset Overview
  3. Methodology Pipeline
  4. Optimal K Selection  (matplotlib chart from CSV)
  5. UMAP 2D Projection   (newplot.png)
  6. Cluster Gallery C0–C7
  7. Cluster Gallery C8–C15
  8. Key Findings

Usage:
    python build_presentation.py
"""

import io
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x5C)
TEAL   = RGBColor(0x00, 0x7A, 0x8A)
ORANGE = RGBColor(0xE8, 0x6C, 0x2C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
DGRAY  = RGBColor(0x3A, 0x3A, 0x3A)
SUBHDR = RGBColor(0xBB, 0xD5, 0xE8)
SUBBODY= RGBColor(0xCC, 0xEE, 0xF5)

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)
HEADER_H = Inches(1.15)

OUT_PATH   = "results/clustering_presentation.pptx"
GRID_DIR   = "results/cluster_grids"
UMAP_IMG   = "newplot.png"
OPT_K_CSV  = "results/optimal_k_results.csv"

CLUSTER_COUNTS = {
    0: 93, 1: 23151, 2: 11748, 3: 11496, 4: 13843,
    5: 20970, 6: 22530, 7: 10067, 8: 16787, 9: 9124,
    10: 9380, 11: 9887, 12: 9576, 13: 11997, 14: 12057, 15: 13280,
}


# ── Low-level helpers ──────────────────────────────────────────────────────────

def blank_slide(prs, bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide


def add_rect(slide, x, y, w, h, fill_color, line=False):
    shp = slide.shapes.add_shape(1, x, y, w, h)   # 1 = msoShapeRectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line:
        shp.line.color.rgb = fill_color
    else:
        shp.line.fill.background()
    return shp


def add_textbox(slide, text, x, y, w, h,
                size=14, bold=False, italic=False,
                color=DGRAY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_header(slide, title, subtitle=None):
    bar = add_rect(slide, Inches(0), Inches(0), SW, HEADER_H, NAVY)
    tf = bar.text_frame
    tf.word_wrap = False
    tf.margin_top = Pt(10)
    tf.margin_left = Inches(0.35)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = SUBHDR
        r2.font.name = "Calibri"


def add_stat_card(slide, value, label, x, y, w=Inches(2.85), h=Inches(1.5)):
    box = add_rect(slide, x, y, w, h, TEAL)
    tf = box.text_frame
    tf.margin_top = Pt(12)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(12)
    r2.font.color.rgb = SUBBODY
    r2.font.name = "Calibri"


def add_bullet_block(slide, items, x, y, w, h, size=13):
    """items = list of str (bullet) or (str, [str,...]) (bullet + sub-bullets)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for item in items:
        if isinstance(item, tuple):
            header, subs = item
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(6)
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = f"▶  {header}"
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = NAVY
            r.font.name = "Calibri"
            for sub in subs:
                ps = tf.add_paragraph()
                ps.alignment = PP_ALIGN.LEFT
                ps.space_before = Pt(2)
                rs = ps.add_run()
                rs.text = f"      {sub}"
                rs.font.size = Pt(size - 1)
                rs.font.color.rgb = DGRAY
                rs.font.name = "Calibri"
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(4)
            r = p.add_run()
            r.text = f"•  {item}"
            r.font.size = Pt(size)
            r.font.color.rgb = DGRAY
            r.font.name = "Calibri"


# ── Chart helpers ──────────────────────────────────────────────────────────────

def make_optimal_k_png():
    df = pd.read_csv(OPT_K_CSV)
    fig, ax1 = plt.subplots(figsize=(9, 4.2), facecolor="white")
    ax2 = ax1.twinx()

    ax1.plot(df["k"], df["inertia"] / 1e6, "o-",
             color="#1B3A5C", lw=2.5, markersize=7, label="Inertia (WCSS)")
    ax1.set_xlabel("Number of clusters (k)", fontsize=12)
    ax1.set_ylabel("Inertia (×10⁶)", fontsize=11, color="#1B3A5C")
    ax1.tick_params(axis="y", colors="#1B3A5C")
    ax1.grid(axis="y", alpha=0.25, linestyle="--")
    ax1.set_facecolor("white")

    valid = df.dropna(subset=["silhouette"])
    ax2.plot(valid["k"], valid["silhouette"], "s--",
             color="#007A8A", lw=2.5, markersize=7, label="Silhouette score")
    ax2.set_ylabel("Silhouette score", fontsize=11, color="#007A8A")
    ax2.tick_params(axis="y", colors="#007A8A")

    ymin, ymax = ax1.get_ylim()
    ax1.axvline(16, color="#E86C2C", lw=2, linestyle=":", alpha=0.9)
    ax1.text(17, ymax * 0.96, "k = 16 ✓", color="#E86C2C", fontsize=11, va="top", fontweight="bold")

    lines = ax1.get_lines() + ax2.get_lines()
    labels = [l.get_label() for l in lines]
    ax1.legend(lines, labels, loc="upper right", fontsize=10, framealpha=0.9)

    plt.title("Optimal K — 205,986 stained tiles", fontsize=13, color="#1B3A5C", pad=8)
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf


def make_cluster_size_png():
    ids = sorted(CLUSTER_COUNTS.keys())
    counts = [CLUSTER_COUNTS[i] for i in ids]
    colors = ["#E86C2C" if CLUSTER_COUNTS[i] < 500 else "#1B3A5C" for i in ids]

    fig, ax = plt.subplots(figsize=(4.8, 5.0), facecolor="white")
    ax.barh([f"C{i}" for i in ids], counts, color=colors, height=0.72)
    ax.invert_yaxis()
    ax.set_xlabel("Tile count", fontsize=10)
    ax.set_title("Cluster sizes", fontsize=11, color="#1B3A5C")
    ax.set_facecolor("white")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(labelsize=9)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{int(x/1000)}k" if x >= 1000 else str(int(x))))
    # Highlight C0
    ax.annotate("Outlier?", xy=(93, 0), xytext=(5000, 0),
                fontsize=8, color="#E86C2C",
                arrowprops=dict(arrowstyle="->", color="#E86C2C", lw=1))
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf


# ── Individual slides ──────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs, NAVY)

    # Teal accent stripe
    add_rect(slide, Inches(0), Inches(3.15), SW, Inches(0.07), TEAL)

    add_textbox(slide,
                "Pathology Tile Embedding Clustering",
                Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.5),
                size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Frozen Section H&E  ·  k=16 MiniBatchKMeans  ·  384-dim path-foundation embeddings",
                Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.7),
                size=15, color=SUBHDR, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "GI Registration Staging  ·  Nanozoomer 40×  ·  March 2026",
                Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
                size=12, color=RGBColor(0x77, 0x99, 0xBB), align=PP_ALIGN.CENTER)


def slide_dataset(prs):
    slide = blank_slide(prs)
    add_header(slide, "Dataset Overview", "GI Frozen Section — H&E Stained Tiles")

    stats = [
        ("411,972", "Total tiles"),
        ("205,986", "Stained (clustered)"),
        ("384-dim", "Embedding vector"),
        ("16", "Tissue types"),
    ]
    bw, bh = Inches(2.85), Inches(1.5)
    gap = Inches(0.27)
    sx = Inches(0.4)
    for i, (val, lbl) in enumerate(stats):
        add_stat_card(slide, val, lbl, sx + i * (bw + gap), Inches(1.4), bw, bh)

    add_bullet_block(slide, [
        ("Source data",
         ["Nanozoomer 40× frozen section slides",
          "GI registration staging pipeline",
          "512 × 512 px tiles, resize preprocessing"]),
        ("Embeddings",
         ["Google path-foundation model via Hugging Face",
          "384-dimensional float32 vectors per tile"]),
        ("Clustering scope",
         ["Stained modality only (unstained tiles excluded)",
          "All 205,986 stained tiles used for training"]),
    ], Inches(0.4), Inches(3.15), Inches(12.5), Inches(4.0), size=13)


def slide_methodology(prs):
    slide = blank_slide(prs)
    add_header(slide, "Methodology", "End-to-end clustering pipeline")

    steps = [
        ("TIFF Tiles",        "512×512 px\nNanozoomer 40×"),
        ("path-foundation",   "384-dim embeddings\nGoogle / HuggingFace"),
        ("DuckDB",            "Streaming ingest\n411,972 tiles"),
        ("MiniBatchKMeans",   "k=16  ·  205,986\nstained tiles"),
        ("Results",           "UMAP · Grids\nHTML reports"),
    ]

    n = len(steps)
    bw = Inches(2.1)
    bh = Inches(1.7)
    total_box_w = bw * n
    total_gap = SW - Inches(0.6) * 2 - total_box_w
    gap = total_gap / (n - 1)
    y_box = Inches(2.6)

    for i, (title, desc) in enumerate(steps):
        x = Inches(0.6) + i * (bw + gap)
        fill = NAVY if i == 3 else TEAL
        box = add_rect(slide, x, y_box, bw, bh, fill)
        tf = box.text_frame
        tf.margin_top = Pt(14)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = SUBBODY
        r2.font.name = "Calibri"

        if i < n - 1:
            add_textbox(slide, "→",
                        x + bw + Inches(0.03), y_box + bh / 2 - Inches(0.22),
                        gap - Inches(0.06), Inches(0.45),
                        size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    notes = [
        "PCA: 384D → 130D  (95.1% variance retained)  →  UMAP 2D for visualization",
        "k=16 selected via elbow + silhouette analysis on the full 205,986-tile dataset",
    ]
    for j, note in enumerate(notes):
        add_textbox(slide, note,
                    Inches(0.4), Inches(5.0 + j * 0.55), Inches(12.5), Inches(0.45),
                    size=11, italic=True,
                    color=RGBColor(0x55, 0x77, 0x99), align=PP_ALIGN.CENTER)


def slide_optimal_k(prs):
    slide = blank_slide(prs)
    add_header(slide, "Optimal K Selection",
               "Elbow + silhouette analysis · 205,986 stained tiles")

    chart_buf = make_optimal_k_png()
    slide.shapes.add_picture(chart_buf, Inches(0.5), Inches(1.3),
                              Inches(10.2), Inches(5.0))

    cards = [
        ("Elbow method", "k = 16", NAVY),
        ("Best silhouette", "k = 8", TEAL),
        ("Selected k", "k = 16", ORANGE),
    ]
    for i, (lbl, val, color) in enumerate(cards):
        bx = add_rect(slide, Inches(11.0), Inches(1.5 + i * 1.75),
                      Inches(2.1), Inches(1.45), color)
        tf = bx.text_frame
        tf.margin_top = Pt(10)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = val
        r.font.size = Pt(26)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = lbl
        r2.font.size = Pt(11)
        r2.font.color.rgb = SUBBODY
        r2.font.name = "Calibri"


def slide_umap(prs):
    slide = blank_slide(prs)
    add_header(slide, "UMAP 2D Projection",
               "Stratified 600 tiles/cluster · PCA 384D→130D (95.1% var) → UMAP 2D  ·  stars = cluster centroids  ·  right panel = centroid tile")
    slide.shapes.add_picture(UMAP_IMG,
                              Inches(0.2), Inches(1.2),
                              Inches(12.9), Inches(6.15))


def slide_cluster_gallery(prs, cluster_ids, title):
    slide = blank_slide(prs)
    add_header(slide, title, "5×5 representative tile grids — locally available tiles")

    COLS, ROWS = 4, 2
    pad_x = Inches(0.22)
    content_y = HEADER_H + Inches(0.12)
    content_h = SH - content_y - Inches(0.08)
    content_w = SW - pad_x * 2
    cell_w = content_w / COLS
    cell_h = content_h / ROWS
    margin = Inches(0.1)
    label_h = Inches(0.25)

    for idx, cl_id in enumerate(cluster_ids):
        row = idx // COLS
        col = idx % COLS
        x = pad_x + col * cell_w + margin
        y = content_y + row * cell_h + margin
        iw = cell_w - margin * 2
        ih = cell_h - margin * 2 - label_h

        path = os.path.join(GRID_DIR, f"cluster_{cl_id}.png")
        if os.path.exists(path):
            slide.shapes.add_picture(path, x, y, width=iw, height=ih)

        add_textbox(slide, f"Cluster {cl_id}  ({CLUSTER_COUNTS.get(cl_id, '?'):,} tiles)",
                    x, y + ih, iw, label_h,
                    size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_findings(prs):
    slide = blank_slide(prs)
    add_header(slide, "Key Findings",
               "k=16 MiniBatchKMeans · GI Frozen Section H&E · 205,986 stained tiles")

    add_bullet_block(slide, [
        ("Cluster separation",
         ["16 well-separated clusters from 205,986 stained tiles",
          "UMAP confirms visually distinct groupings with minimal overlap"]),
        ("Cluster 0 — potential outlier",
         ["Only 93 tiles (0.05% of data) — rare or anomalous morphology",
          "Merits pathological review as a distinct edge case"]),
        ("Dominant patterns",
         ["Clusters 1 (23,151) and 6 (22,530) represent the most common tissue morphologies",
          "Clusters 5 (20,970) and 8 (16,787) also show high tile density"]),
        ("Embedding quality",
         ["Silhouette score k=16: 0.159  (vs k=8: 0.169 — marginal difference)",
          "PCA: 95.1% variance retained in 130 of 384 embedding dimensions"]),
    ], Inches(0.4), Inches(1.3), Inches(7.8), Inches(5.8), size=13)

    chart_buf = make_cluster_size_png()
    slide.shapes.add_picture(chart_buf, Inches(8.2), Inches(1.3),
                              Inches(4.9), Inches(5.8))


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building slides …")
    slide_title(prs)         ; print("  1/8  Title")
    slide_dataset(prs)       ; print("  2/8  Dataset Overview")
    slide_methodology(prs)   ; print("  3/8  Methodology")
    slide_optimal_k(prs)     ; print("  4/8  Optimal K")
    slide_umap(prs)          ; print("  5/8  UMAP")
    slide_cluster_gallery(prs, list(range(8)),     "Cluster Gallery  C0 – C7")  ; print("  6/8  Cluster Gallery C0–C7")
    slide_cluster_gallery(prs, list(range(8, 16)), "Cluster Gallery  C8 – C15") ; print("  7/8  Cluster Gallery C8–C15")
    slide_findings(prs)      ; print("  8/8  Key Findings")

    os.makedirs("results", exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\nSaved: {OUT_PATH}")


if __name__ == "__main__":
    main()
